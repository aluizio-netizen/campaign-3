# -*- coding: utf-8 -*-
"""
Gera a apresentacao em PowerPoint para a aula "Somos um em Cristo"
Escola Biblica Dominical - Segunda Igreja Batista de Campo Grande-MS
Base: comentario de Romanos 11-12 (Caps. 18 e 19 do material anexo).
Aula unica de 60 minutos. Tipografia ampliada para boa leitura na projecao.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------- Paleta ----------------
AZUL      = RGBColor(0x1F, 0x3A, 0x5F)   # azul profundo
AZUL_ESC  = RGBColor(0x16, 0x2A, 0x45)
DOURADO   = RGBColor(0xC9, 0xA2, 0x27)   # dourado
CREME     = RGBColor(0xF6, 0xF2, 0xE8)   # creme claro
BRANCO    = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_TX  = RGBColor(0x2B, 0x2B, 0x2B)
TEAL      = RGBColor(0x2E, 0x6E, 0x6A)

FONTE_T = "Georgia"        # titulos
FONTE_C = "Calibri"        # corpo

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def add_slide(bg=CREME):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    sp = r._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def faixa(s, top, height, color):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, EMU_W, height)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


def caixa(s, left, top, width, height):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def setp(p, text, size, color, bold=False, font=FONTE_C, italic=False,
         align=PP_ALIGN.LEFT, space_after=8, space_before=0):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    r = p.runs[0]
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = color
    return p


def bullet(tf, text, size=26, color=CINZA_TX, bold=False, level=0,
           first=False, space_after=14, marker="• ", italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = (marker + text) if marker else text
    p.level = level
    p.space_after = Pt(space_after)
    r = p.runs[0]
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONTE_C
    r.font.color.rgb = color; r.font.italic = italic
    return p


def rodape(s, num, dark=False):
    cor = CREME if dark else AZUL
    tb, tf = caixa(s, Inches(0.4), Inches(7.0), Inches(9), Inches(0.4))
    setp(tf.paragraphs[0], "Somos um em Cristo  ·  EBD – 2ª IB Campo Grande-MS",
         12, cor, font=FONTE_C)
    tb2, tf2 = caixa(s, Inches(12.3), Inches(7.0), Inches(0.9), Inches(0.4))
    setp(tf2.paragraphs[0], str(num), 12, cor, align=PP_ALIGN.RIGHT)


def titulo_secao(s, kicker, titulo):
    """Slide divisor de secao (fundo azul)."""
    faixa(s, Inches(3.0), Inches(0.07), DOURADO)
    tb, tf = caixa(s, Inches(1.0), Inches(1.9), Inches(11.3), Inches(1.0))
    setp(tf.paragraphs[0], kicker, 28, DOURADO, bold=True,
         font=FONTE_C, align=PP_ALIGN.CENTER)
    tb2, tf2 = caixa(s, Inches(0.7), Inches(3.15), Inches(11.9), Inches(2.0))
    setp(tf2.paragraphs[0], titulo, 50, BRANCO, bold=True,
         font=FONTE_T, align=PP_ALIGN.CENTER)


def cabecalho(s, titulo, ref=None):
    """Cabecalho padrao de slide de conteudo."""
    faixa(s, 0, Inches(1.3), AZUL)
    faixa(s, Inches(1.3), Inches(0.09), DOURADO)
    tb, tf = caixa(s, Inches(0.6), Inches(0.22), Inches(12.1), Inches(0.95))
    setp(tf.paragraphs[0], titulo, 34, BRANCO, bold=True, font=FONTE_T)
    if ref:
        tb2, tf2 = caixa(s, Inches(0.62), Inches(0.95), Inches(12.1), Inches(0.4))
        setp(tf2.paragraphs[0], ref, 18, DOURADO, bold=True, italic=True)


def versiculo_box(s, top, texto, ref, height=Inches(1.7), txt=24, ref_sz=19):
    """Caixa destacada para versiculo biblico."""
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), top,
                              Inches(11.5), height)
    card.fill.solid(); card.fill.fore_color.rgb = BRANCO
    card.line.color.rgb = DOURADO; card.line.width = Pt(1.5)
    card.shadow.inherit = False
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), top,
                             Inches(0.15), height)
    bar.fill.solid(); bar.fill.fore_color.rgb = DOURADO
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = card.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.45); tf.margin_right = Inches(0.4)
    setp(tf.paragraphs[0], "“" + texto + "”", txt, AZUL_ESC,
         italic=True, font=FONTE_T, space_after=4)
    p2 = tf.add_paragraph()
    setp(p2, ref, ref_sz, TEAL, bold=True, align=PP_ALIGN.RIGHT)


def conteudo(titulo, ref, bullets, num, top=Inches(1.75)):
    """Slide padrao de bullets grandes."""
    s = add_slide()
    cabecalho(s, titulo, ref)
    tb, tf = caixa(s, Inches(0.9), top, Inches(11.6), Inches(5.0))
    for i, (txt, kw) in enumerate(bullets):
        bullet(tf, txt, first=(i == 0), **kw)
    rodape(s, num)
    return s


# =====================================================================
# SLIDE 1 - CAPA
# =====================================================================
s = add_slide(AZUL)
faixa(s, 0, EMU_H, AZUL)
faixa(s, Inches(2.5), Inches(0.07), DOURADO)
faixa(s, Inches(5.45), Inches(0.07), DOURADO)
tb, tf = caixa(s, Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.6))
setp(tf.paragraphs[0], "ESCOLA BÍBLICA DOMINICAL  ·  CLASSE DE ADULTOS",
     18, DOURADO, bold=True, align=PP_ALIGN.CENTER)
tb, tf = caixa(s, Inches(0.6), Inches(2.65), Inches(12.1), Inches(2.0))
setp(tf.paragraphs[0], "SOMOS UM EM CRISTO", 64, BRANCO, bold=True,
     font=FONTE_T, align=PP_ALIGN.CENTER)
tb, tf = caixa(s, Inches(0.8), Inches(4.55), Inches(11.7), Inches(0.8))
setp(tf.paragraphs[0], "Uma só oliveira, um só corpo, um só amor",
     26, CREME, italic=True, font=FONTE_T, align=PP_ALIGN.CENTER)
tb, tf = caixa(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.6))
setp(tf.paragraphs[0], "Estudo em Romanos 11 e 12", 22, DOURADO,
     bold=True, align=PP_ALIGN.CENTER)
tb, tf = caixa(s, Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.5))
setp(tf.paragraphs[0],
     "Segunda Igreja Batista de Campo Grande-MS  ·  Duração: 60 minutos",
     16, CREME, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 2 - VERSICULO TEMA
# =====================================================================
s = add_slide(AZUL_ESC)
faixa(s, 0, EMU_H, AZUL_ESC)
tb, tf = caixa(s, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.6))
setp(tf.paragraphs[0], "VERSÍCULO-TEMA", 22, DOURADO, bold=True,
     align=PP_ALIGN.CENTER)
tb, tf = caixa(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.2))
setp(tf.paragraphs[0],
     "“Assim nós, conquanto muitos, somos um só corpo em Cristo, "
     "e membros uns dos outros.”",
     40, BRANCO, italic=True, font=FONTE_T, align=PP_ALIGN.CENTER, space_after=10)
tb, tf = caixa(s, Inches(1.2), Inches(5.5), Inches(10.9), Inches(0.7))
setp(tf.paragraphs[0], "Romanos 12.5", 28, DOURADO, bold=True,
     align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 3 - ROTEIRO
# =====================================================================
conteudo("Para onde vamos hoje", None, [
    ("Em Cristo, Deus forma UM só povo de muitos povos.",
     dict(size=27, color=AZUL, bold=True)),
    ("1.  Um só povo — a oliveira de Deus (Romanos 11)",
     dict(size=25, color=CINZA_TX)),
    ("2.  Um só corpo — unidade, diversidade, mutualidade e utilidade "
     "(Romanos 12.1-8)", dict(size=25, color=CINZA_TX)),
    ("3.  Um só amor — a vida que sustenta a unidade (Romanos 12.9-21)",
     dict(size=25, color=CINZA_TX)),
    ("Aplicação: como vivemos essa unidade na nossa igreja?",
     dict(size=25, color=TEAL, bold=True)),
], 3)

# =====================================================================
# SLIDE 4 - INTRODUCAO
# =====================================================================
conteudo("Introdução: da doutrina à vida", "Romanos 1–11 → 12–16", [
    ("Romanos é o maior tratado teológico do NT. Nos capítulos 1–11, "
     "Paulo expõe a doutrina; a partir do 12, trata da vida prática.",
     dict(size=26, color=CINZA_TX)),
    ("O capítulo 11 termina em adoração: “A Ele seja a glória para "
     "sempre!” (Rm 11.36). A unidade nasce da adoração.",
     dict(size=26, color=CINZA_TX)),
    ("Não há relacionamento vertical correto com Deus se os "
     "relacionamentos horizontais com os irmãos estão errados.",
     dict(size=26, color=AZUL, bold=True)),
], 4)

# =====================================================================
# SLIDE 5 - DIVISOR PONTO 1
# =====================================================================
s = add_slide(AZUL)
faixa(s, 0, EMU_H, AZUL)
titulo_secao(s, "PONTO 1", "Um só povo: a oliveira de Deus")
tb, tf = caixa(s, Inches(1.0), Inches(5.1), Inches(11.3), Inches(0.8))
setp(tf.paragraphs[0], "Romanos 11", 24, CREME, italic=True,
     align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 6 - DEUS NAO REJEITOU
# =====================================================================
conteudo("Deus não rejeitou o seu povo", "Romanos 11.1-10", [
    ("“De modo nenhum!” Deus não rejeitou Israel. O próprio Paulo é a "
     "prova: israelita, da descendência de Abraão.",
     dict(size=26, color=CINZA_TX)),
    ("A rejeição é apenas parcial: sempre houve um remanescente fiel — "
     "como os 7.000 que não dobraram os joelhos a Baal.",
     dict(size=26, color=CINZA_TX)),
    ("Esse remanescente existe “segundo a eleição da graça” (11.5). "
     "Não por obras — senão a graça já não seria graça.",
     dict(size=26, color=AZUL, bold=True)),
], 6)

# =====================================================================
# SLIDE 7 - GENTIOS ENXERTADOS
# =====================================================================
conteudo("Gentios enxertados na oliveira", "Romanos 11.11-24", [
    ("A queda de Israel abriu espaço para a salvação dos gentios — e "
     "isso deve provocar Israel a um santo ciúme, para a sua "
     "restauração.", dict(size=25, color=CINZA_TX)),
    ("A raiz é Abraão; os ramos naturais são os judeus; os gentios são "
     "ramos silvestres enxertados pela fé.", dict(size=25, color=CINZA_TX)),
    ("“Não te glories contra os ramos” (11.18): o gentio não sustenta a "
     "raiz — a raiz é que o sustenta. Humildade, não orgulho.",
     dict(size=25, color=AZUL, bold=True)),
], 7)

# =====================================================================
# SLIDE 8 - A VERDADE CENTRAL (oliveira)
# =====================================================================
s = add_slide(TEAL)
faixa(s, 0, EMU_H, TEAL)
tb, tf = caixa(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.9))
setp(tf.paragraphs[0], "A verdade central", 26, CREME, bold=True,
     align=PP_ALIGN.CENTER)
versiculo_box(s, Inches(1.85),
    "Uma só oliveira representa todos os salvos, sem importar a origem. "
    "Para judeus e gentios, a salvação é a mesma: pela graça, "
    "por meio da fé, sobre a base da expiação de Cristo.",
    "Comentário de Romanos 11", height=Inches(2.3), txt=25)
tb, tf = caixa(s, Inches(0.9), Inches(4.6), Inches(11.5), Inches(2.4))
bullet(tf, "Judeus e gentios foram chamados a formar “um só corpo em "
        "Cristo, a saber, a igreja” (Ef 3.4-6).", 26, BRANCO, bold=True,
        first=True)
bullet(tf, "A unidade do povo de Deus não é uniformidade — é união na "
        "mesma raiz, pela mesma graça.", 25, CREME)

# =====================================================================
# SLIDE 9 - DIVISOR PONTO 2
# =====================================================================
s = add_slide(AZUL)
faixa(s, 0, EMU_H, AZUL)
titulo_secao(s, "PONTO 2", "Um só corpo em Cristo")
tb, tf = caixa(s, Inches(1.0), Inches(5.1), Inches(11.3), Inches(0.8))
setp(tf.paragraphs[0], "Romanos 12.1-8", 24, CREME, italic=True,
     align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 10 - VIDAS TRANSFORMADAS
# =====================================================================
s = add_slide()
cabecalho(s, "A base: vidas transformadas", "Romanos 12.1-2")
tb, tf = caixa(s, Inches(0.9), Inches(1.7), Inches(11.6), Inches(2.4))
bullet(tf, "Primeiro, o relacionamento com Deus: oferecer o corpo em "
        "sacrifício vivo, santo e agradável — nosso culto racional.",
        26, CINZA_TX, first=True)
bullet(tf, "Não se conformar com este mundo, mas transformar-se pela "
        "renovação da mente (gr.: metamorfose, de dentro para fora).",
        26, CINZA_TX)
versiculo_box(s, Inches(4.9),
    "Apresentai o vosso corpo por sacrifício vivo... "
    "transformai-vos pela renovação da vossa mente.",
    "Romanos 12.1-2", height=Inches(1.7), txt=23)
rodape(s, 10)

# =====================================================================
# SLIDE 11 - QUATRO VERDADES (cards)
# =====================================================================
s = add_slide()
cabecalho(s, "Um só corpo: quatro verdades", "Romanos 12.3-8")
labels = [
    ("UNIDADE", "Somos um só corpo\nRm 12.5", AZUL),
    ("DIVERSIDADE", "Muitos membros,\ndons diferentes\nRm 12.4-6", TEAL),
    ("MUTUALIDADE", "Membros uns\ndos outros\nRm 12.5", DOURADO),
    ("UTILIDADE", "Dons para edificar\no corpo\nRm 12.6-8", AZUL_ESC),
]
x = Inches(0.7)
for titulo, desc, cor in labels:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.1),
                              Inches(2.85), Inches(3.7))
    card.fill.solid(); card.fill.fore_color.rgb = cor
    card.line.fill.background(); card.shadow.inherit = False
    tf = card.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    txtcor = AZUL_ESC if cor == DOURADO else BRANCO
    setp(tf.paragraphs[0], titulo, 23, txtcor,
         bold=True, font=FONTE_T, align=PP_ALIGN.CENTER, space_after=12)
    for ln in desc.split("\n"):
        p = tf.add_paragraph()
        setp(p, ln, 18, (AZUL_ESC if cor == DOURADO else CREME),
             align=PP_ALIGN.CENTER, space_after=2)
    x = Emu(int(x) + int(Inches(3.05)))
rodape(s, 11)

# =====================================================================
# SLIDES 12-14 - UNIDADE / DIVERSIDADE / MUTUALIDADE
# =====================================================================
conteudo("1. Unidade", "Romanos 12.5", [
    ("“Somos um só corpo.” Uma só família, um só rebanho.",
     dict(size=28, color=AZUL, bold=True)),
    ("O que dá unidade ao corpo é estarmos ligados à mesma Cabeça — "
     "Cristo — e irrigados pelo mesmo sangue.", dict(size=26, color=CINZA_TX)),
    ("A unidade não é algo que produzimos; é algo que recebemos. Já "
     "somos um em Cristo — cabe-nos guardá-la.", dict(size=26, color=CINZA_TX)),
], 12)

conteudo("2. Diversidade", "Romanos 12.4-6a", [
    ("“Nem todos os membros têm a mesma função.”",
     dict(size=28, color=TEAL, bold=True)),
    ("Pessoas das mais diversas origens, ambientes e temperamentos são "
     "dotadas por Deus de grande variedade de dons.",
     dict(size=26, color=CINZA_TX)),
    ("A diversidade não ameaça a unidade — ela a enriquece. Cada um "
     "coopera para o bem de todos.", dict(size=26, color=CINZA_TX)),
], 13)

conteudo("3. Mutualidade", "Romanos 12.5", [
    ("“Somos membros uns dos outros.”",
     dict(size=28, color=DOURADO, bold=True)),
    ("Não estamos competindo — estamos servindo uns aos outros, com "
     "igual cuidado uns pelos outros (1Co 12.25).",
     dict(size=26, color=CINZA_TX)),
    ("Precisamos uns dos outros. Ninguém é dispensável; ninguém é "
     "autossuficiente.", dict(size=26, color=CINZA_TX)),
], 14)

# =====================================================================
# SLIDE 15 - UTILIDADE / DONS
# =====================================================================
s = add_slide()
cabecalho(s, "4. Utilidade: dons para servir", "Romanos 12.6-8")
tb, tf = caixa(s, Inches(0.9), Inches(1.6), Inches(11.6), Inches(0.9))
bullet(tf, "Os dons são dados pela Trindade para edificar o corpo — "
        "nenhum para exibição própria:", 26, AZUL, bold=True, first=True,
        space_after=8)
dons = [
    "Profecia — proclamar a Palavra",
    "Ministério / serviço — servir com amor",
    "Ensino — instruir no entendimento",
    "Exortação — encorajar e consolar",
    "Contribuição — repartir com liberalidade",
    "Presidência / liderança — dirigir com diligência",
    "Misericórdia — acolher com alegria",
]
tb, tf = caixa(s, Inches(1.1), Inches(2.65), Inches(11.4), Inches(4.2))
for i, d in enumerate(dons):
    bullet(tf, d, 23, CINZA_TX, first=(i == 0), space_after=8)
rodape(s, 15)

# =====================================================================
# SLIDE 16 - DIVISOR PONTO 3
# =====================================================================
s = add_slide(AZUL)
faixa(s, 0, EMU_H, AZUL)
titulo_secao(s, "PONTO 3", "Um só amor")
tb, tf = caixa(s, Inches(1.0), Inches(5.1), Inches(11.3), Inches(0.8))
setp(tf.paragraphs[0], "Romanos 12.9-21", 24, CREME, italic=True,
     align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 17 - AMOR SINCERO (versiculo + intro)
# =====================================================================
s = add_slide()
cabecalho(s, "O amor que sustenta a unidade", "Romanos 12.9-10")
versiculo_box(s, Inches(1.75),
    "O amor seja sem hipocrisia. Detestai o mal, apegando-vos ao bem.",
    "Romanos 12.9", height=Inches(1.6), txt=25)
tb, tf = caixa(s, Inches(0.9), Inches(3.7), Inches(11.6), Inches(3.1))
bullet(tf, "O amor (ágape) é o sistema circulatório do corpo "
        "espiritual: faz todos os membros funcionarem de forma saudável.",
        26, AZUL, bold=True, first=True)
bullet(tf, "Amor sincero, sem máscara, com discernimento: detestar o "
        "mal, apegar-se ao bem.", 25, CINZA_TX)
bullet(tf, "Afeição fraternal e honra mútua: “preferindo-vos em honra "
        "uns aos outros” (12.10).", 25, CINZA_TX)
rodape(s, 17)

# =====================================================================
# SLIDE 18 - QUATRO PORTAS ABERTAS
# =====================================================================
s = add_slide()
cabecalho(s, "Um amor que se abre", "Romanos 12.13-16")
items = [
    ("Coração aberto", "amor fraternal, sincero e cordial (12.9-10)"),
    ("Lábios abertos", "abençoar e não amaldiçoar (12.14)"),
    ("Mãos abertas", "repartir com os santos nas necessidades (12.13a)"),
    ("Casa aberta", "praticar a hospitalidade — acolher (12.13b)"),
]
y = Inches(1.75)
for titulo, desc in items:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y,
                              Inches(11.5), Inches(1.12))
    card.fill.solid(); card.fill.fore_color.rgb = CREME
    card.line.color.rgb = DOURADO; card.line.width = Pt(1.25)
    card.shadow.inherit = False
    tf = card.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.4)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = titulo + "  —  "
    r1.font.bold = True; r1.font.size = Pt(24); r1.font.name = FONTE_C
    r1.font.color.rgb = AZUL
    r2 = p.add_run(); r2.text = desc
    r2.font.size = Pt(22); r2.font.name = FONTE_C; r2.font.color.rgb = CINZA_TX
    y = Emu(int(y) + int(Inches(1.28)))
rodape(s, 18)

# =====================================================================
# SLIDE 19 - VENCER O MAL COM O BEM
# =====================================================================
s = add_slide()
cabecalho(s, "Vencer o mal com o bem", "Romanos 12.17-21")
tb, tf = caixa(s, Inches(0.9), Inches(1.7), Inches(11.6), Inches(2.6))
bullet(tf, "A unidade cristã alcança até os inimigos: não retaliar, "
        "não vingar-se, não guardar mágoa.", 26, CINZA_TX, first=True)
bullet(tf, "“Se possível, quanto depender de vós, tende paz com todos "
        "os homens” (12.18). O cristão é pacificador.",
        26, CINZA_TX)
versiculo_box(s, Inches(5.0),
    "Não te deixes vencer do mal, mas vence o mal com o bem.",
    "Romanos 12.21", height=Inches(1.5), txt=25)
rodape(s, 19)

# =====================================================================
# SLIDE 20 - APLICACAO
# =====================================================================
s = add_slide(AZUL)
faixa(s, 0, EMU_H, AZUL)
faixa(s, Inches(1.3), Inches(0.09), DOURADO)
tb, tf = caixa(s, Inches(0.6), Inches(0.32), Inches(12), Inches(0.95))
setp(tf.paragraphs[0], "E na nossa igreja?", 34, BRANCO, bold=True,
     font=FONTE_T)
tb, tf = caixa(s, Inches(0.9), Inches(1.7), Inches(11.6), Inches(5.1))
bullet(tf, "Guardo a unidade ou alimento divisões e fofocas? (12.16)",
       25, CREME, first=True)
bullet(tf, "Conheço e uso o meu dom para edificar o corpo — ou vivo como "
        "espectador?", 25, CREME)
bullet(tf, "Valorizo a diversidade de dons e pessoas, ou exijo que todos "
        "sejam iguais a mim?", 25, CREME)
bullet(tf, "Meu amor é sincero? Minha casa, mãos e lábios estão abertos "
        "aos irmãos?", 25, CREME)
bullet(tf, "Há alguém com quem preciso fazer as pazes esta semana?",
       25, DOURADO, bold=True)

# =====================================================================
# SLIDE 21 - CONCLUSAO
# =====================================================================
s = add_slide()
cabecalho(s, "Conclusão")
tb, tf = caixa(s, Inches(0.9), Inches(1.75), Inches(11.6), Inches(3.2))
bullet(tf, "Em Cristo, Deus fez de muitos um só povo: uma só oliveira, "
        "um só corpo, um só amor.", 28, AZUL, bold=True, first=True)
bullet(tf, "A unidade é dom recebido pela graça; a diversidade nos "
        "enriquece; o amor mantém tudo unido.", 26, CINZA_TX)
bullet(tf, "Ser um em Cristo não é um ideal distante: é uma realidade a "
        "ser vivida, todos os dias, uns com os outros.",
        26, CINZA_TX)
versiculo_box(s, Inches(5.6),
    "A Ele seja a glória para sempre! Amém.",
    "Romanos 11.36", height=Inches(1.2), txt=25)
rodape(s, 21)

# =====================================================================
# SLIDE 22 - ENCERRAMENTO
# =====================================================================
s = add_slide(AZUL_ESC)
faixa(s, 0, EMU_H, AZUL_ESC)
faixa(s, Inches(2.8), Inches(0.07), DOURADO)
tb, tf = caixa(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(1.4))
setp(tf.paragraphs[0], "SOMOS UM EM CRISTO", 52, BRANCO, bold=True,
     font=FONTE_T, align=PP_ALIGN.CENTER)
tb, tf = caixa(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.9))
setp(tf.paragraphs[0],
     "“Conquanto muitos, somos um só corpo em Cristo” — Rm 12.5",
     24, DOURADO, italic=True, align=PP_ALIGN.CENTER)
tb, tf = caixa(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.6))
setp(tf.paragraphs[0],
     "Escola Bíblica Dominical · Segunda Igreja Batista de Campo Grande-MS",
     16, CREME, align=PP_ALIGN.CENTER)

prs.save("/home/user/campaign-3/aula-somos-um-em-cristo/Somos_um_em_Cristo_apresentacao.pptx")
print("PPTX gerado com", len(prs.slides._sldIdLst), "slides.")
